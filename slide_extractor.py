import os
import sys
import shutil
import tkinter as tk
from tkinter import filedialog, ttk, messagebox
import subprocess
from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageTk
import glob

def resource_path(relative_path):
    """ Get absolute path to resource (used for PyInstaller binary access) """
    try:
        base_path = sys._MEIPASS  # PyInstaller uses this
    except AttributeError:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)

def open_path(path):
    """Open a file or directory with the default OS application."""
    try:
        os.startfile(path)
    except AttributeError:  # Not on Windows
        opener = "open" if sys.platform == "darwin" else "xdg-open"
        subprocess.run([opener, path])

def show_preview(directory):
    """Display the first extracted slide as a thumbnail."""
    images = sorted(glob.glob(os.path.join(directory, "*.jpg")))
    if not images:
        preview_label.config(image="")
        preview_label.image = None
        return
    img = Image.open(images[0])
    img.thumbnail((200, 200))
    img_tk = ImageTk.PhotoImage(img)
    preview_label.config(image=img_tk)
    preview_label.image = img_tk
    img.close()

def select_video():
    filepath = filedialog.askopenfilename(
        title="Select Video",
        filetypes=[("Video files", "*.mp4 *.mov *.avi *.mkv")]
    )
    if filepath:
        video_path.set(filepath)
        basename = os.path.splitext(os.path.basename(filepath))[0]
        out_dir_var.set(f"{basename}_slides")
        status.set("Ready")

def extract_slides():
    path = video_path.get()
    if not path:
        status.set("Please select a video file.")
        return

    try:
        threshold = float(threshold_val.get())
    except ValueError:
        status.set("Invalid threshold.")
        return

    directory = out_dir_var.get().strip()
    if not directory:
        basename = os.path.splitext(os.path.basename(path))[0]
        directory = f"{basename}_slides"
    os.makedirs(directory, exist_ok=True)
    out_dir_var.set(directory)

    ffmpeg_path = ffmpeg_path_var.get()
    if not os.path.isfile(ffmpeg_path):
        ffmpeg_path = shutil.which("ffmpeg")
        if not ffmpeg_path:
            status.set("ffmpeg not found. Set path in Settings.")
            return
    timestamps = [t.strip() for t in timestamps_var.get().split(',') if t.strip()]
    output_pattern = os.path.join(directory, "%04d.jpg")
    cmds = []
    if timestamps:
        for i, ts in enumerate(timestamps, start=1):
            out_file = os.path.join(directory, f"{i:04d}.jpg")
            cmds.append([ffmpeg_path, "-ss", ts, "-i", path, "-frames:v", "1", out_file])
    else:
        cmds.append([
            ffmpeg_path, "-i", path,
            "-filter_complex", f"select=gt(scene\\,{threshold})",
            "-vsync", "vfr", output_pattern
        ])

    try:
        status.set("Extracting...")
        progress['value'] = 0
        if timestamps:
            progress.configure(mode='determinate', maximum=len(cmds))
        else:
            progress.configure(mode='indeterminate')
            progress.start()
        root.update()
        for idx, cmd in enumerate(cmds, start=1):
            subprocess.run(cmd, check=True)
            if timestamps:
                progress['value'] = idx
                root.update_idletasks()
        progress.stop()
        status.set(f"Done! Saved to: {directory}")
        if open_after_var.get():
            open_path(directory)
        show_preview(directory)
    except Exception as e:
        progress.stop()
        status.set(f"Error: {e}")

# Export functions
def export_to_pptx(directory):
    if not directory or not os.path.isdir(directory):
        status.set("No slides to export.")
        return
    images = sorted(glob.glob(os.path.join(directory, "*.jpg")))
    if not images:
        status.set("No slides to export.")
        return
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for img in images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            img,
            Inches(0),
            Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )
    basename = os.path.basename(directory)
    pptx_path = os.path.join(directory, f"{basename}.pptx")
    try:
        prs.save(pptx_path)
        status.set(f"Exported to {pptx_path}")
        if open_after_var.get():
            open_path(pptx_path)
        if delete_var.get():
            for img in images:
                os.remove(img)
            show_preview(directory)
    except Exception as e:
        status.set(f"Error: {e}")


def export_to_pdf(directory):
    if not directory or not os.path.isdir(directory):
        status.set("No slides to export.")
        return
    images = sorted(glob.glob(os.path.join(directory, "*.jpg")))
    if not images:
        status.set("No slides to export.")
        return
    img_objs = [Image.open(p).convert("RGB") for p in images]
    basename = os.path.basename(directory)
    pdf_path = os.path.join(directory, f"{basename}.pdf")
    try:
        img_objs[0].save(pdf_path, save_all=True, append_images=img_objs[1:])
        status.set(f"Exported to {pdf_path}")
        if open_after_var.get():
            open_path(pdf_path)
        if delete_var.get():
            for img_path in images:
                os.remove(img_path)
            show_preview(directory)
    except Exception as e:
        status.set(f"Error: {e}")
    finally:
        for img in img_objs:
            img.close()

# GUI Setup
root = tk.Tk()
root.title("Slide Extractor")

video_path = tk.StringVar()
threshold_val = tk.StringVar(value="0.2")
out_dir_var = tk.StringVar()
timestamps_var = tk.StringVar()
status = tk.StringVar(value="Select a video.")
delete_var = tk.BooleanVar()
open_after_var = tk.BooleanVar(value=True)
ffmpeg_path_var = tk.StringVar(value=resource_path("ffmpeg.exe"))
dark_mode_var = tk.BooleanVar()

style = ttk.Style()

def apply_theme():
    """Toggle light or dark appearance."""
    if dark_mode_var.get():
        style.theme_use("clam")
        style.configure("TFrame", background="#2e2e2e")
        style.configure("TLabel", background="#2e2e2e", foreground="white")
        style.configure("TEntry", fieldbackground="#4d4d4d", foreground="white")
        style.configure("TButton", background="#4d4d4d")
        style.configure("TCheckbutton", background="#2e2e2e", foreground="white")
        root.configure(bg="#2e2e2e")
    else:
        style.theme_use("clam")
        style.configure("TFrame", background="")
        style.configure("TLabel", background="", foreground="")
        style.configure("TEntry", fieldbackground="white", foreground="black")
        style.configure("TButton", background="")
        style.configure("TCheckbutton", background="", foreground="")
        root.configure(bg="")

frm = ttk.Frame(root, padding=10)
frm.pack(fill=tk.BOTH, expand=True)

ttk.Label(frm, text="Video Path:").grid(row=0, column=0, sticky="w")
ttk.Entry(frm, textvariable=video_path, width=50).grid(row=0, column=1, sticky="ew")
ttk.Button(frm, text="Browse", command=select_video).grid(row=0, column=2, padx=5)

ttk.Label(frm, text="Output Directory:").grid(row=1, column=0, sticky="w")
ttk.Entry(frm, textvariable=out_dir_var, width=50).grid(row=1, column=1, sticky="ew")
ttk.Button(frm, text="Browse", command=lambda: out_dir_var.set(filedialog.askdirectory() or out_dir_var.get())).grid(row=1, column=2, padx=5)

ttk.Label(frm, text="Scene Threshold (e.g. 0.2):").grid(row=2, column=0, sticky="w")
ttk.Entry(frm, textvariable=threshold_val, width=10).grid(row=2, column=1, sticky="w")

ttk.Label(frm, text="Timestamps (comma separated, optional):").grid(row=3, column=0, sticky="w")
ttk.Entry(frm, textvariable=timestamps_var, width=50).grid(row=3, column=1, columnspan=2, sticky="ew")

progress = ttk.Progressbar(frm, mode="determinate")
progress.grid(row=4, column=0, columnspan=3, pady=5, sticky="ew")

preview_label = ttk.Label(frm)
preview_label.grid(row=5, column=0, columnspan=3, pady=5)

ttk.Button(frm, text="Extract Slides", command=extract_slides).grid(row=6, column=0, pady=5)
ttk.Button(frm, text="Export to PPTX", command=lambda: export_to_pptx(out_dir_var.get())).grid(row=6, column=1)
ttk.Button(frm, text="Export to PDF", command=lambda: export_to_pdf(out_dir_var.get())).grid(row=6, column=2)

ttk.Label(frm, textvariable=status, foreground="blue").grid(row=7, column=0, columnspan=3, pady=5)

frm.columnconfigure(1, weight=1)

menubar = tk.Menu(root)
settings_menu = tk.Menu(menubar, tearoff=0)
settings_menu.add_checkbutton(label="Open files when done", variable=open_after_var)
settings_menu.add_checkbutton(label="Delete JPGs after export", variable=delete_var)
settings_menu.add_checkbutton(label="Dark mode", variable=dark_mode_var, command=apply_theme)
def set_ffmpeg_path():
    path = filedialog.askopenfilename(title="Select ffmpeg", filetypes=[("ffmpeg", "ffmpeg*"), ("All files", "*")])
    if path:
        ffmpeg_path_var.set(path)
settings_menu.add_command(label="Set ffmpeg path", command=set_ffmpeg_path)
menubar.add_cascade(label="Settings", menu=settings_menu)
usage_menu = tk.Menu(menubar, tearoff=0)
def show_usage():
    message = (
        "1. Select a video file.\n"
        "2. Adjust the threshold or add timestamps.\n"
        "3. Choose an output folder and click Extract.\n"
        "4. Export to PPTX or PDF if desired."
    )
    messagebox.showinfo("Usage", message)
usage_menu.add_command(label="How to Use", command=show_usage)
menubar.add_cascade(label="Usage", menu=usage_menu)
root.config(menu=menubar)

apply_theme()

root.mainloop()
